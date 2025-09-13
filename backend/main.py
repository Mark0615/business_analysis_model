import re, time, uuid, logging
import concurrent.futures as futures
import numpy as np
import pandas as pd
import io, json, os
import matplotlib 
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from typing import List, Optional, Dict, Any, Tuple
from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import JSONResponse, PlainTextResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

load_dotenv()

OPENAI_MODEL = os.getenv("DECK_MODEL", "gpt-4o-mini")
DECK_MAX_SLIDES = int(os.getenv("DECK_MAX_SLIDES", "10"))
DECK_STRICT_LLM = os.getenv("DECK_STRICT_LLM", "1").lower() in ("1","true","yes")

api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise RuntimeError("Missing OPENAI_API_KEY. Put it in backend/.env or set an env var.")
client = OpenAI(api_key=api_key)

CORS_ORIGINS_RAW = os.getenv(
    "CORS_ORIGINS",
    "http://localhost:3000,http://127.0.0.1:3000"
)
_raw_origins = [o.strip() for o in CORS_ORIGINS_RAW.split(",") if o.strip()]
_plain_origins = [o for o in _raw_origins if "*" not in o]
_wildcards = [o for o in _raw_origins if "*" in o]
if _wildcards:
    import re as _re
    _patterns = [_re.escape(w).replace(r"\*", r"[^/]+") for w in _wildcards]
    CORS_ORIGIN_REGEX = r"^(?:%s)$" % "|".join(_patterns)
else:
    CORS_ORIGIN_REGEX = None

REQUEST_MAX_BYTES = int(os.getenv("REQUEST_MAX_BYTES", "10485760"))
MAX_ROWS = int(os.getenv("MAX_ROWS", "200000"))
ALGO_TIMEOUT_SEC = int(os.getenv("ALGO_TIMEOUT_SEC", "20"))
LOG_LEVEL = os.getenv("LOG_LEVEL", "INFO").upper()
REQUIRE_HTTPS = os.getenv("REQUIRE_HTTPS", "0").lower() in ("1", "true", "yes")
SENTRY_DSN = os.getenv("SENTRY_DSN", "")

logging.basicConfig(
    level=getattr(logging, LOG_LEVEL, logging.INFO),
    format="%(asctime)s %(levelname)s %(message)s",
)
if SENTRY_DSN:
    try:
        import sentry_sdk
        sentry_sdk.init(dsn=SENTRY_DSN, traces_sample_rate=0.05)
        logging.info("Sentry initialized")
    except Exception as e:
        logging.warning(f"Sentry init failed: {e}")

def _lazy_import_ml():
    from sklearn.cluster import KMeans
    from mlxtend.preprocessing import TransactionEncoder
    from mlxtend.frequent_patterns import apriori, association_rules
    return KMeans, TransactionEncoder, apriori, association_rules

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=_plain_origins,
    allow_origin_regex=CORS_ORIGIN_REGEX,
    allow_credentials=True,
    allow_methods=["POST", "GET", "OPTIONS"],
    allow_headers=["*"],
)

START_TS = time.time()

@app.middleware("http")
async def request_context(request: Request, call_next):
    rid = request.headers.get("x-request-id") or str(uuid.uuid4())
    request.state.rid = rid
    t0 = time.time()
    try:
        resp = await call_next(request)
    except Exception:
        logging.exception(f"[{rid}] unhandled error")
        return JSONResponse({"error": "internal_server_error", "rid": rid}, status_code=500)
    dt_ms = int((time.time() - t0) * 1000)
    resp.headers["X-Request-ID"] = rid
    logging.info(f"[{rid}] {request.method} {request.url.path} -> {resp.status_code} in {dt_ms}ms")
    return resp

@app.middleware("http")
async def https_enforcer(request: Request, call_next):
    if REQUIRE_HTTPS:
        proto = request.headers.get("x-forwarded-proto") or request.url.scheme
        if proto != "https":
            return JSONResponse({"error": "https_required"}, status_code=400)
    return await call_next(request)

@app.middleware("http")
async def security_headers(request: Request, call_next):
    resp = await call_next(request)
    resp.headers["X-Content-Type-Options"] = "nosniff"
    resp.headers["X-Frame-Options"] = "DENY"
    resp.headers["Referrer-Policy"] = "no-referrer"
    if REQUIRE_HTTPS:
        resp.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains; preload"
    return resp

@app.middleware("http")
async def size_guard(request: Request, call_next):
    if request.method in ("POST", "PUT", "PATCH"):
        cl = request.headers.get("content-length")
        if cl and int(cl) > REQUEST_MAX_BYTES:
            return JSONResponse({"error": "request_too_large", "limit_bytes": REQUEST_MAX_BYTES}, status_code=413)
        if not cl:
            body = await request.body()
            if len(body) > REQUEST_MAX_BYTES:
                return JSONResponse({"error": "request_too_large", "limit_bytes": REQUEST_MAX_BYTES}, status_code=413)
    return await call_next(request)

class Row(BaseModel):
    date: str
    product_name: str
    product_category: Optional[str] = None
    quantity: Optional[float] = 0
    revenue: float
    customer_name: Optional[str] = None
    customer_email: Optional[str] = None
    order_id: Optional[str] = None

class AnalyzeIn(BaseModel):
    rows: List[Row]

class AnalyzeAssocIn(BaseModel):
    rows: List[Row]
    date_from: str
    date_to: str

class DeckIn(BaseModel):
    title: Optional[str] = "Business Product Analysis"
    insights: Dict[str, Any]

def _pick_for_prompt(ins: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "kpis": {"mom": ins.get("mom"), "wow": ins.get("wow")},
        "monthly": (ins.get("monthly") or [])[-12:],
        "top5_products": (ins.get("top5_products") or [])[:5],
        "rfm_share": ins.get("rfm_share"),
        "rfm_counts": ins.get("rfm_counts"),
        "rfm_avg_monetary": ins.get("rfm_avg_monetary"),
        "assoc_rules": (ins.get("assoc_rules") or [])[:3],
        "potential_products": (ins.get("potential_products") or [])[:5],
        "nurture_mid": (ins.get("nurture_mid") or [])[:10],
        "nurture_low": (ins.get("nurture_low") or [])[:10],
    }

def _ask_llm_for_outline(title: str, payload: Dict[str, Any]) -> List[Dict[str, Any]]:
    system = (
        "你是商業營運顧問與簡報撰寫者，使用繁體中文，語氣專業且精簡。"
        "根據提供的分析 JSON，產出 6~12 張投影片的大綱。"
        "每張投影片需包含 title 與 3~5 個 bullet（每個 bullet 20 字以內）。"
        "包含：封面、整體摘要、趨勢、暢銷商品、RFM 客群洞察、關聯規則（若有）、潛力商品、行動建議與後續實施。"
        "只輸出 JSON 物件：{\"slides\": [{\"title\": str, \"bullets\": [str, ...]}]}。不得輸出多餘文字。"
    )
    user = {
        "deck_title": title,
        "max_slides": DECK_MAX_SLIDES,
        "insights": _pick_for_prompt(payload),
    }

    schema = {
        "name": "deck_outline",
        "schema": {
            "type": "object",
            "properties": {
                "slides": {
                    "type": "array",
                    "minItems": 6,
                    "maxItems": DECK_MAX_SLIDES,
                    "items": {
                        "type": "object",
                        "required": ["title", "bullets"],
                        "properties": {
                            "title":   {"type": "string", "minLength": 2, "maxLength": 60},
                            "bullets": {"type": "array", "minItems": 3, "maxItems": 5,
                                        "items": {"type": "string", "minLength": 2, "maxLength": 40}}
                        }
                    }
                }
            },
            "required": ["slides"],
            "additionalProperties": False
        },
        "strict": True,
    }

    msgs = [
        {"role": "system", "content": system},
        {"role": "user", "content": json.dumps(user, ensure_ascii=False)},
    ]

    raw = ""
    try:
        # 優先: SDK（responses + json_schema）
        try:
            resp = client.responses.create(
                model=OPENAI_MODEL,
                temperature=0.2,
                max_output_tokens=1200,
                response_format={"type": "json_schema", "json_schema": schema},
                input=msgs,
            )
            raw = resp.output_text or ""
        except TypeError:
            # 舊 SDK（沒有 response_format）：退回純文字
            resp = client.responses.create(
                model=OPENAI_MODEL,
                temperature=0.2,
                input=msgs,
            )
            raw = (resp.output_text or "").strip()
        except AttributeError:
            # 更舊（沒有 responses API）：退到 chat.completions
            comp = client.chat.completions.create(
                model=OPENAI_MODEL,
                temperature=0.2,
                messages=msgs,
            )
            raw = (comp.choices[0].message.content or "").strip()

        s = raw.strip()
        if not s.startswith("{"):
            i = s.find("{")
            if i >= 0:
                s = s[i:]
        if "}" in s:
            s = s[: s.rfind("}") + 1]

        obj = json.loads(s)
        slides = obj.get("slides") or obj
        cleaned: List[Dict[str, Any]] = []
        if isinstance(slides, list):
            for item in slides[:DECK_MAX_SLIDES]:
                if not isinstance(item, dict): continue
                t = str(item.get("title") or "").strip()[:60]
                bs = [str(b)[:40] for b in (item.get("bullets") or [])][:5]
                if t and bs: cleaned.append({"title": t, "bullets": bs})
        if cleaned:
            return cleaned
        raise ValueError("model did not return valid slides")
    except Exception as e:
        logging.warning("LLM outline failed, fallback used: %s", e)

    return [
        {"title": title, "bullets": ["資料摘要", "趨勢與商品重點", "後續行動建議"]},
        {"title": "關鍵洞察摘要", "bullets": ["MoM/WoW 變化", "Top 產品", "RFM 概況"]},
    ]

def _build_trend_image(monthly: List[Dict[str, Any]]) -> io.BytesIO:
    buf = io.BytesIO()
    if not monthly:
        return buf
    try:
        xs = [m["yyyymm"] for m in monthly]
        ys = [m["revenue"] for m in monthly]
        plt.figure(figsize=(6, 3.2), dpi=160)
        plt.plot(xs, ys, linewidth=2)
        plt.fill_between(xs, ys, alpha=0.15)
        plt.xticks(rotation=45, ha="right", fontsize=8)
        plt.ylabel("Revenue")
        plt.tight_layout()
        plt.savefig(buf, format="png")
    finally:
        plt.close('all')
    buf.seek(0)
    return buf

def _add_trend_slide(prs, monthly):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "每月營收趨勢（近一年）"
    try:
        img = _build_trend_image(monthly)
        if img.getbuffer().nbytes > 0:
            slide.shapes.add_picture(img, Inches(0.6), Inches(1.6), width=Inches(9.2))
    except Exception:
        pass
    return slide

def _build_top5_chart(top5: List[Dict[str, Any]]) -> io.BytesIO:
    buf = io.BytesIO()
    if not top5: return buf
    names = [x["product_name"] for x in top5]
    vals  = [x["revenue"] for x in top5]
    plt.figure(figsize=(6,3.2), dpi=160)
    plt.barh(names[::-1], vals[::-1])
    plt.tight_layout()
    plt.savefig(buf, format="png"); plt.close(); buf.seek(0)
    return buf

def _build_rfm_donut(share: Dict[str, float]) -> io.BytesIO:
    buf = io.BytesIO()
    if not share: return buf
    labels = ["高價值","中價值","低價值"]
    sizes  = [share.get("high",0), share.get("mid",0), share.get("low",0)]
    plt.figure(figsize=(6,3.2), dpi=160)
    wedges, _ = plt.pie(sizes, startangle=140, wedgeprops=dict(width=0.45))
    plt.legend(wedges, labels, loc="center left", bbox_to_anchor=(1,0.5))
    plt.tight_layout()
    plt.savefig(buf, format="png"); plt.close(); buf.seek(0)
    return buf

def _add_bullet_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for b in bullets[:5]:
        p = tf.add_paragraph()
        p.text = str(b)
        p.level = 0
    return slide

def _add_cover(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    st = slide.shapes.title.text_frame.paragraphs[0]
    st.font.size = Pt(40)
    st.alignment = PP_ALIGN.LEFT
    sub = slide.placeholders[1]
    sub.text = "Automated deck"
    return slide

def _add_image_slide(prs, title: str, img_buf: io.BytesIO):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if img_buf.getbuffer().nbytes > 0:
        slide.shapes.add_picture(img_buf, Inches(0.6), Inches(1.6), width=Inches(9.2))
    return slide

def _add_assoc_slide(prs, rules: List[Dict[str, Any]]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "關聯規則（Top）"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    show = rules[:5] if rules else []
    if not show:
        p = tf.add_paragraph(); p.text = "暫無顯著規則"; p.level = 0
        return slide
    for r in show:
        p = tf.add_paragraph()
        p.text = f"{' + '.join(r['antecedents'])} → {' + '.join(r['consequents'])}｜Lift {r['lift']:.2f}｜Conf {(r['confidence']*100):.1f}%"
        p.level = 0
    return slide


@app.get("/ping")
def ping():
    return {"status": "ok", "uptime_sec": int(time.time() - START_TS)}

@app.get("/healthz", response_class=PlainTextResponse)
def healthz():
    return "ok"

def _assoc_rules(df: pd.DataFrame) -> Tuple[list, Dict[str, Any]]:
    try:
        KMeans, TransactionEncoder, apriori, association_rules = _lazy_import_ml()
    except Exception:
        return [], {"ok": False, "baskets_total": 0, "baskets_valid": 0, "reason": "algorithm_error"}

    cust_key = df.get("customer_email", pd.Series([""] * len(df))).fillna("").str.strip()
    if not (cust_key.str.len() > 0).any():
        cust_key = df.get("customer_name", pd.Series([""] * len(df))).fillna("").str.strip()

    if df.get("order_id", pd.Series([None] * len(df))).notna().any():
        df["basket_id"] = df["order_id"].fillna("NA").astype(str)
    else:
        df["basket_id"] = cust_key + "_" + df["date"].dt.strftime("%Y-%m-%d")

    tx_series = df.groupby("basket_id")["product_name"].apply(lambda s: sorted(set(s)))
    tx_all = int(len(tx_series))
    tx = [t for t in tx_series.tolist() if len(t) >= 2]

    assoc_rules_out = []
    reason = None
    try:
        if len(tx) >= 5:
            te = TransactionEncoder()
            tx_df = pd.DataFrame(te.fit(tx).transform(tx), columns=te.columns_)
            freq = apriori(tx_df, min_support=0.005, use_colnames=True)
            if not freq.empty:
                rules = association_rules(freq, metric="lift", min_threshold=1.05)
                if not rules.empty:
                    rules = rules.sort_values(["lift", "confidence"], ascending=False).head(5)
                    for _, r in rules.iterrows():
                        assoc_rules_out.append({
                            "antecedents": sorted(list(r["antecedents"])),
                            "consequents": sorted(list(r["consequents"])),
                            "support": round(float(r["support"]), 4),
                            "confidence": round(float(r["confidence"]), 4),
                            "lift": round(float(r["lift"]), 4),
                        })
                else:
                    reason = "no_significant_rules"
            else:
                reason = "too_sparse_support"
        else:
            reason = "single_item_baskets_or_too_few"
    except Exception:
        reason = "algorithm_error"

    diag = {
        "ok": len(assoc_rules_out) > 0,
        "baskets_total": tx_all,
        "baskets_valid": int(len(tx)),
        "reason": reason,
    }
    return assoc_rules_out, diag

def _rfm(df: pd.DataFrame) -> Tuple[Dict[str, float], Dict[str, int], Dict[str, float], Dict[str, list], Dict[str, Any]]:
    try:
        KMeans, *_ = _lazy_import_ml()
    except Exception:
        KMeans = None

    cust_key = df.get("customer_email", pd.Series([""] * len(df))).fillna("").str.strip()
    if not (cust_key.str.len() > 0).any():
        cust_key = df.get("customer_name", pd.Series([""] * len(df))).fillna("").str.strip()
    df = df.join(cust_key.rename("customer_id"))

    if df.get("order_id", pd.Series([None] * len(df))).notna().any():
        df["basket_id"] = df["order_id"].fillna("NA").astype(str)
    else:
        df["basket_id"] = df["customer_id"].fillna("") + "_" + df["date"].dt.strftime("%Y-%m-%d")
    df["order_key"] = df["basket_id"]

    cust = (df.groupby("customer_id")
              .agg(last_date=("date", "max"),
                   orders=("order_key", pd.Series.nunique),
                   monetary=("revenue", "sum"))
              .reset_index())
    cust = cust[cust["customer_id"].fillna("").str.len() > 0]

    customers_total = int(cust["customer_id"].nunique())
    customers_used = int(len(cust))

    rfm_share = {"low": 0.0, "mid": 0.0, "high": 0.0}
    rfm_counts = {"low": 0, "mid": 0, "high": 0}
    rfm_avg_monetary = {"low": 0.0, "mid": 0.0, "high": 0.0}
    nurture: Dict[str, list] = {"mid": [], "low": []}
    reason = None

    if len(cust) >= 3:
        max_date = df["date"].max()
        cust["recency"] = (max_date - cust["last_date"]).dt.days

        use_quantile = False
        if KMeans is not None:
            try:
                X = cust[["recency", "orders", "monetary"]].astype(float)
                X = (X - X.mean()) / (X.std(ddof=0) + 1e-9)
                km = KMeans(n_clusters=3, n_init=10, random_state=42)
                cust["segment"] = km.fit_predict(X)

                counts = cust["segment"].value_counts()
                if (len(counts) < 3) or (counts.min() <= 1):
                    use_quantile = True
                else:
                    rank = (cust.groupby("segment")["monetary"].mean().sort_values().reset_index())
                    rank["level"] = ["low", "mid", "high"]
                    m = dict(zip(rank["segment"], rank["level"]))
                    cust["level"] = cust["segment"].map(m)
            except Exception:
                use_quantile = True
        else:
            use_quantile = True

        if use_quantile:
            q1, q2 = cust["monetary"].quantile([1/3, 2/3])
            def lvl(v):
                if v <= q1: return "low"
                elif v <= q2: return "mid"
                else: return "high"
            cust["level"] = cust["monetary"].apply(lvl)

        levels = ["low", "mid", "high"]

        cnt_map = cust["level"].value_counts().to_dict()
        rfm_counts = {lv: int(cnt_map.get(lv, 0)) for lv in levels}

        avg_map = {lv: float(cust.loc[cust["level"] == lv, "monetary"].mean() or 0.0) for lv in levels}
        rfm_avg_monetary = {lv: round(avg_map.get(lv, 0.0), 2) for lv in levels}

        total_rev = float(cust["monetary"].sum() or 0.0)
        share_sum_map = {lv: float(cust.loc[cust["level"] == lv, "monetary"].sum()) for lv in levels}
        rfm_share = {lv: (round(share_sum_map.get(lv, 0.0) / total_rev, 4) if total_rev > 0 else 0.0) for lv in levels}

        r_th = 45
        m_th_mid = cust.loc[cust["level"] == "mid", "monetary"].median() if (cust["level"] == "mid").any() else 0.0
        m_th_low = cust.loc[cust["level"] == "low", "monetary"].median() if (cust["level"] == "low").any() else 0.0

        def build_list(df_seg, m_th):
            seg = df_seg[(df_seg["recency"] > r_th) & (df_seg["monetary"] >= m_th)]
            out = (seg.sort_values("monetary", ascending=False)
                    .head(200)[["customer_id", "last_date", "orders", "monetary"]]
                    .rename(columns={
                        "customer_id": "customer_email",
                        "last_date": "last_order_date",
                        "orders": "lifetime_orders",
                        "monetary": "lifetime_revenue"
                    })
                    .to_dict(orient="records"))
            return out

        nurture["mid"] = build_list(cust[cust["level"] == "mid"], m_th_mid)
        nurture["low"] = build_list(cust[cust["level"] == "low"], m_th_low)
    else:
        reason = "no_customer_key_or_too_few"

    diag = {
        "ok": customers_used >= 3,
        "customers_total": customers_total,
        "customers_used": customers_used,
        "reason": reason
    }
    return rfm_share, rfm_counts, rfm_avg_monetary, nurture, diag

def _potential_products(df: pd.DataFrame,
                        top5_names: set,
                        monthly_ordered: pd.DataFrame) -> list:
    out = []
    df["yyyymm"] = df["date"].dt.to_period("M").astype(str)
    prod_m = (df.groupby(["product_name", "yyyymm"], as_index=False)["revenue"].sum()
                .sort_values(["product_name", "yyyymm"]))
    if not prod_m.empty and len(monthly_ordered) >= 3:
        last3 = sorted(monthly_ordered["yyyymm"].unique())[-3:]
        recent = prod_m[prod_m["yyyymm"].isin(last3)]
        if not recent.empty:
            g = recent.pivot(index="product_name", columns="yyyymm", values="revenue").fillna(0.0)
            if g.shape[1] == 3:
                g["mom_last"] = (g.iloc[:, 2] - g.iloc[:, 1]) / (g.iloc[:, 1].replace(0, np.nan))
                g["sum3"] = g.sum(axis=1)
                cand = g[(g["mom_last"] > 0.2) & (~g.index.isin(top5_names))]
                cand = cand.sort_values(["mom_last", "sum3"], ascending=False).head(5)
                for name, row in cand.iterrows():
                    out.append({
                        "product_name": name,
                        "mom_last": round(float(row["mom_last"]), 3),
                        "sum3_revenue": round(float(row["sum3"]), 2)
                    })
    return out

# ---------- endpoints ----------
@app.post("/generate_deck",
    responses={
        200: {
            "content": {
                "application/vnd.openxmlformats-officedocument.presentationml.presentation": {}
            },
            "description": "Generate PPTX file."
        }
    },
)
def generate_deck(in_: DeckIn):
    try:
        outline = _ask_llm_for_outline(in_.title or "Business Product Analysis", in_.insights)

        prs = Presentation()
        _add_cover(prs, in_.title or "Business Product Analysis")

        monthly = (in_.insights.get("monthly") or [])[-12:]
        if monthly:
            _add_trend_slide(prs, monthly)

        top5 = in_.insights.get("top5_products") or []
        if top5:
            _add_image_slide(prs, "暢銷商品 Top 5（營收）", _build_top5_chart(top5))

        rfm_share = in_.insights.get("rfm_share") or {}
        if rfm_share:
            _add_image_slide(prs, "RFM 客群佔比", _build_rfm_donut(rfm_share))

        rules = in_.insights.get("assoc_rules") or []
        _add_assoc_slide(prs, rules)

        if DECK_STRICT_LLM and (not outline or not isinstance(outline, list)):
            raise HTTPException(status_code=502, detail={"error":"llm_invalid_output","message":"model did not return valid slides"})

        for s in outline or []:
            _add_bullet_slide(prs, s.get("title") or "Untitled", s.get("bullets") or [])

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        fn = f"business_deck_{int(time.time())}.pptx"
        return StreamingResponse(
            out,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{fn}"', "X-Deck-Source": "llm"}
        )
    except HTTPException:
        raise
    except Exception as e:
        logging.exception("generate_deck failed")
        raise HTTPException(status_code=502, detail={"error": "deck_failed", "message": str(e)[:300]})

@app.post("/assoc_window")
def assoc_window(payload: AnalyzeAssocIn, request: Request):
    n = len(payload.rows or [])
    if n == 0:
        return {"error": "no_rows",
                "diagnostics": {"assoc": {"ok": False, "reason": "empty_payload"}}}

    df = pd.DataFrame([r.dict() for r in payload.rows])
    df["date"] = pd.to_datetime(df["date"], errors="coerce")
    df["product_name"] = df["product_name"].astype(str).str.strip()
    df = df.dropna(subset=["date", "product_name"])

    try:
        s = pd.to_datetime(payload.date_from)
        e = pd.to_datetime(payload.date_to)
    except Exception:
        raise HTTPException(status_code=422, detail={"error": "invalid_dates"})

    win_df = df[(df["date"] >= s) & (df["date"] <= e)].copy()
    if win_df.empty:
        return {
            "assoc_rules": [],
            "diagnostics": {
                "assoc": {"ok": False, "reason": "window_empty", "baskets_total": 0, "baskets_valid": 0},
                "window": {"from": payload.date_from, "to": payload.date_to, "rows": 0}
            }
        }

    assoc_rules, assoc_diag = _assoc_rules(win_df)
    return {
        "assoc_rules": assoc_rules,
        "diagnostics": {
            "assoc": assoc_diag,
            "window": {"from": payload.date_from, "to": payload.date_to, "rows": int(len(win_df))}
        }
    }

@app.post("/analyze")
async def analyze(payload: AnalyzeIn, request: Request):
    rid = getattr(request.state, "rid", "-")

    n = len(payload.rows or [])
    if n == 0:
        return {"error": "no_rows",
                "diagnostics": {"required": {"ok": False, "reason": "empty_payload"}}}
    if n > MAX_ROWS:
        raise HTTPException(status_code=422, detail={"error": "too_many_rows", "max_rows": MAX_ROWS, "got": n})

    df = pd.DataFrame([r.dict() for r in payload.rows])

    required_cols = ["date", "product_name", "revenue"]
    missing = [c for c in required_cols if c not in df.columns]
    if missing:
        return {"error": "missing_required_fields",
                "diagnostics": {"required": {"ok": False, "reason": "missing_fields", "missing": missing}}}

    df["date"] = pd.to_datetime(df["date"], errors="coerce")
    df["product_name"] = df["product_name"].astype(str).str.strip()
    df["quantity"] = pd.to_numeric(df.get("quantity", 0), errors="coerce").fillna(0)
    df["revenue"] = pd.to_numeric(df["revenue"], errors="coerce").fillna(0.0)
    df = df.dropna(subset=["date", "product_name"])
    if df.empty:
        return {"error": "invalid_required_values",
                "diagnostics": {"required": {"ok": False, "reason": "all_rows_invalid"}}}

    top5_df = (df.groupby("product_name", as_index=False)["revenue"].sum()
                 .sort_values("revenue", ascending=False).head(5))
    top5 = top5_df.to_dict(orient="records")

    monthly_df = (df.assign(yyyymm=df["date"].dt.strftime("%Y-%m"))
                    .groupby("yyyymm", as_index=False)["revenue"].sum()
                    .sort_values("yyyymm"))

    mom = None
    if len(monthly_df) >= 2:
        a, b = monthly_df.iloc[-2]["revenue"], monthly_df.iloc[-1]["revenue"]
        mom = None if a == 0 else (b - a) / a

    last_day = df["date"].max()
    r1 = df.loc[(df["date"] >= last_day - pd.Timedelta(days=13)) & (df["date"] < last_day - pd.Timedelta(days=6)), "revenue"].sum()
    r2 = df.loc[(df["date"] >= last_day - pd.Timedelta(days=6)) & (df["date"] <= last_day), "revenue"].sum()
    wow = None if r1 == 0 else (r2 - r1) / r1

    assoc_rules, assoc_diag = [], {"ok": False, "reason": "timeout", "baskets_total": 0, "baskets_valid": 0}
    rfm_share = {"low": 0.0, "mid": 0.0, "high": 0.0}
    rfm_counts = {"low": 0, "mid": 0, "high": 0}
    rfm_avg_monetary = {"low": 0.0, "mid": 0.0, "high": 0.0}
    nurture_mid: list = []
    nurture_low: list = []
    rfm_diag = {"ok": False, "reason": "timeout", "customers_total": 0, "customers_used": 0}

    with futures.ThreadPoolExecutor(max_workers=2) as ex:
        f_assoc = ex.submit(_assoc_rules, df.copy())
        f_rfm   = ex.submit(_rfm, df.copy())

        try:
            assoc_rules, assoc_diag = f_assoc.result(timeout=ALGO_TIMEOUT_SEC)
        except futures.TimeoutError:
            assoc_diag = {"ok": False, "reason": "timeout", "baskets_total": 0, "baskets_valid": 0}
            logging.warning(f"[{rid}] assoc timeout")

        try:
            rfm_share, rfm_counts, rfm_avg_monetary, nurture_dict, rfm_diag = f_rfm.result(timeout=ALGO_TIMEOUT_SEC)
            nurture_mid = nurture_dict.get("mid", [])
            nurture_low = nurture_dict.get("low", [])
        except futures.TimeoutError:
            rfm_diag = {"ok": False, "reason": "timeout", "customers_total": 0, "customers_used": 0}
            logging.warning(f"[{rid}] rfm timeout")

    potential_products = _potential_products(df.copy(), set(top5_df["product_name"].tolist()), monthly_df)

    return {
        "top5_products": top5,
        "monthly": monthly_df.to_dict(orient="records"),
        "mom": mom,
        "wow": wow,
        "assoc_rules": assoc_rules,
        "rfm_share": rfm_share,
        "rfm_counts": rfm_counts,
        "rfm_avg_monetary": rfm_avg_monetary,
        "nurture_list": (nurture_mid + nurture_low),
        "nurture_mid": nurture_mid,
        "nurture_low": nurture_low,
        "potential_products": potential_products,
        "diagnostics": {"assoc": assoc_diag, "rfm": rfm_diag, "required": {"ok": True}},
    }
